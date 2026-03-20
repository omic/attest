"""Browser-based multi-LLM chat — uses your existing ChatGPT, Claude, and Gemini subscriptions.

Launch: attest chat [--mode browser] [--providers chatgpt,claude,gemini]

Opens a Chromium browser with persistent login sessions. On first run, you'll
log into each chat app once — after that, sessions are remembered. Your question
is sent to all providers simultaneously via their web UIs.

Judging/scoring/synthesis uses a free API tier (Gemini or Groq) so you don't
need any API keys for the primary chat providers.

Requires: pip install playwright && playwright install chromium
"""

from __future__ import annotations

import asyncio
import json
import logging
import os
import re
import sys
import time
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Readline tab-completion for @ (providers, files) and / (commands)
# ---------------------------------------------------------------------------

_COMMANDS = [
    "/tournament", "/share", "/consensus", "/new", "/model", "/auto",
    "/web", "/screenshot", "/export", "/upload", "/output", "/downloads",
    "/providers", "/save", "/load", "/artifacts", "/help", "/quit", "/exit",
]

_PROVIDER_NAMES = ["chatgpt", "claude", "gemini"]


class _ChatCompleter:
    """Tab-completion for the attest chat prompt.

    Supports:
    - /cmd     → command completion
    - @prov    → @chatgpt, @claude, @gemini (selective send)
    - @file    → files and directories in cwd
    - @dir/    → files inside that directory
    """

    def __init__(self, cwd: str):
        self.cwd = cwd
        self._matches: list[str] = []

    def update_cwd(self, cwd: str):
        self.cwd = cwd

    def complete(self, text: str, state: int) -> str | None:
        if state == 0:
            self._matches = self._get_matches(text)
        return self._matches[state] if state < len(self._matches) else None

    def _get_matches(self, text: str) -> list[str]:
        """Build completion list based on current token."""
        try:
            import readline
            line = readline.get_line_buffer()
            begin = readline.get_begidx()

            # If cursor is right after @, readline gives us the text after @
            if begin > 0 and line[begin - 1:begin] == "@":
                return self._complete_at(text)
        except (ImportError, AttributeError):
            pass

        # Fallback: infer from the text itself
        if text.startswith("@"):
            return self._complete_at(text[1:], prefix="@")
        if text.startswith("/"):
            return self._complete_command(text)

        return []

    def _complete_command(self, text: str) -> list[str]:
        return [c + " " for c in _COMMANDS if c.startswith(text)]

    def _complete_at(self, partial: str, prefix: str = "") -> list[str]:
        """Complete after @.

        First offers provider names, then file/directory matches.
        """
        matches = []

        # Provider names (for selective sending)
        for p in _PROVIDER_NAMES:
            if p.startswith(partial.lower()):
                matches.append(f"{prefix}{p} ")

        # File and directory completion
        if "/" in partial:
            # Completing inside a directory: @docs/rea → @docs/readme.md
            dir_part, file_part = partial.rsplit("/", 1)
            search_dir = os.path.join(self.cwd, dir_part)
        else:
            dir_part = ""
            file_part = partial
            search_dir = self.cwd

        try:
            entries = os.listdir(search_dir)
        except OSError:
            entries = []

        for entry in sorted(entries):
            if entry.startswith("."):
                continue  # Skip hidden files
            if entry.lower().startswith(file_part.lower()):
                full_path = os.path.join(search_dir, entry)
                if dir_part:
                    display = f"{dir_part}/{entry}"
                else:
                    display = entry

                if os.path.isdir(full_path):
                    matches.append(f"{prefix}{display}/")
                else:
                    matches.append(f"{prefix}{display} ")

        return matches


def _setup_readline(cwd: str) -> _ChatCompleter:
    """Configure readline for tab-completion in the chat prompt."""
    try:
        import readline
    except ImportError:
        return _ChatCompleter(cwd)

    completer = _ChatCompleter(cwd)

    readline.set_completer(completer.complete)
    # Treat @ and / as word-break characters so completion triggers properly
    readline.set_completer_delims(" \t\n")
    readline.parse_and_bind("tab: complete")

    return completer


# Terminal colors
_COLORS = {
    "chatgpt": "\033[32m",     # green
    "claude": "\033[95m",      # light magenta
    "gemini": "\033[34m",      # blue
    "consensus": "\033[97;42m",  # white on green
    "system": "\033[90m",      # gray
    "judge": "\033[33m",       # yellow
}
_RESET = "\033[0m"
_BOLD = "\033[1m"

PROFILE_DIR = Path.home() / ".attestdb" / "browser-profile"

MAX_FILE_CHARS = 10_000
FILE_REF_PATTERN = re.compile(r"@([\w./\-*\[\]{}]+)")

# Stall detection: if a provider's response matches these, it's asking for files
# instead of reviewing the content we already sent.
_STALL_PATTERNS = re.compile(
    r"(upload|send.*(files?|documents?)|share.*(files?|documents?)|"
    r"attach|drag.*(drop|and)|didn.t come through|no files|"
    r"can.t (access|see|find|open).*(?:files?|directory|folder)|"
    r"don.t have (access|the files)|paste the contents|"
    r"re-upload|provide the files|I.ll need)",
    re.IGNORECASE,
)
_STALL_MAX_LEN = 1500  # Short responses with stall patterns = likely stuck

# Follow-up detection: provider is asking for human input
_FOLLOWUP_PATTERNS = re.compile(
    r"(would you like|do you want|can you|could you|shall I|should I|"
    r"which (one|option|approach|file)|please (clarify|specify|confirm|choose|select)|"
    r"before I (proceed|continue|start|begin)|let me know|"
    r"a few questions|some questions|clarify.*\?|more (details|info|context))",
    re.IGNORECASE,
)
# Minimum question marks at end of response to trigger follow-up
_MIN_QUESTIONS = 2

# Prompt wrapper when files are included
_FILE_CONTEXT_PREFIX = (
    "IMPORTANT: The complete contents of all relevant files from my working "
    "directory are included below in this message. Do NOT ask me to upload, "
    "attach, share, or paste files — everything is already here. "
    "Review the file contents directly.\n\n"
)

# Auto-response when a provider asks for files instead of reviewing
_STALL_NUDGE = (
    "The complete file contents are already included in my previous message above. "
    "Please scroll up and review the text between the --- markers. "
    "Each file's content is shown inline. Do not ask me to upload anything — "
    "just review what's there and give me your analysis."
)

# CSS selectors for chat UIs.
# These WILL break when UIs change — keep them here for easy updates.
# Each provider has multiple fallback selectors for response extraction.
PROVIDER_SELECTORS = {
    "chatgpt": {
        "url": "https://chatgpt.com",
        "input": 'div[id="prompt-textarea"]',
        "input_fallbacks": [
            'div[contenteditable="true"][id="prompt-textarea"]',
            'div[contenteditable="true"]',
        ],
        "send": 'button[data-testid="send-button"]',
        "send_fallbacks": [
            'button[aria-label="Send prompt"]',
            'button[aria-label="Send"]',
        ],
        "response": 'div[data-message-author-role="assistant"]',
        "response_fallbacks": [
            'article[data-testid^="conversation-turn"] div.markdown',
            'div[data-message-author-role="assistant"] div.markdown',
            'div[data-message-author-role="assistant"]',
            # Newer ChatGPT DOM variants
            'div.agent-turn div.markdown',
            'div.markdown.prose',
            '[data-message-author-role="assistant"] .markdown',
        ],
        "streaming_indicator": 'button[aria-label="Stop generating"]',
        "streaming_fallbacks": [
            'button[aria-label="Stop streaming"]',
            'button[aria-label="Stop"]',
            'div[data-message-author-role="assistant"][data-is-streaming="true"]',
            # Newer: the stop button changes label
            'button[data-testid="stop-button"]',
        ],
        "new_chat": 'a[data-testid="create-new-chat-button"]',
        # ChatGPT canvas / code blocks
        "artifact_selectors": [
            'div[class*="code-block"] pre code',
            'pre code',
            'div[data-testid="canvas-content"]',
        ],
    },
    "gemini": {
        "url": "https://gemini.google.com/app",
        "input": 'div.ql-editor',
        "input_fallbacks": [
            'div[contenteditable="true"]',
            'rich-textarea div[contenteditable="true"]',
        ],
        "send": 'button[aria-label="Send message"]',
        "send_fallbacks": [
            'button.send-button',
            'button[data-test-id="send-button"]',
        ],
        "response": 'message-content',
        "response_fallbacks": [
            'message-content .markdown',
            'message-content',
            'model-response message-content',
            '.response-container-content',
            'div.response-content',
        ],
        "streaming_indicator": 'button[aria-label="Stop response"]',
        "streaming_fallbacks": [
            'mat-icon[data-mat-icon-name="stop_circle"]',
        ],
        "artifact_selectors": [
            'code-block pre code',
            'pre code',
        ],
    },
    "claude": {
        "url": "https://claude.ai/new",
        "input": 'div[contenteditable="true"]',
        "input_fallbacks": [
            'fieldset div[contenteditable="true"]',
            'div.ProseMirror[contenteditable="true"]',
        ],
        "send": 'button[aria-label="Send Message"]',
        "send_fallbacks": [
            'button[aria-label="Send message"]',
            'button[type="button"][class*="send"]',
        ],
        "response": 'div.font-claude-message',
        "response_fallbacks": [
            'div.font-claude-message',
            'div[data-is-streaming] div.grid',
            'div[class*="message"] div[class*="markdown"]',
            'div[data-testid="assistant-message"]',
        ],
        "streaming_indicator": 'button[aria-label="Stop Response"]',
        "streaming_fallbacks": [
            'button[aria-label="Stop response"]',
            'div[data-is-streaming="true"]',
        ],
        # Claude artifacts appear in a separate panel
        "artifact_selectors": [
            'div[data-testid="artifact-content"]',
            'div[class*="artifact"] pre code',
            'div[class*="code-block"] pre code',
            'pre code',
        ],
    },
}

# ---------------------------------------------------------------------------
# Model / tier selectors for each provider's web UI
# ---------------------------------------------------------------------------

# Known models per provider — keys are what the user types, values are the
# display text we look for in the model picker dropdown.
PROVIDER_MODELS = {
    "chatgpt": {
        "4o": "GPT-4o",
        "4o-mini": "GPT-4o mini",
        "o1": "o1",
        "o3": "o3",
        "o4-mini": "o4-mini",
        "deep-research": "Deep Research",
    },
    "claude": {
        "haiku": "Haiku",
        "sonnet": "Sonnet",
        "opus": "Opus",
    },
    "gemini": {
        "flash": "Flash",
        "pro": "Pro",
        "deep-research": "Deep Research",
    },
}

# CSS selectors for the model picker button and dropdown items.
# These WILL break when UIs change — keep them here for easy updates.
MODEL_SELECTORS = {
    "chatgpt": {
        # The model picker is a button/dropdown at the top of the chat
        "picker_button": [
            'button[data-testid="model-switcher-dropdown-button"]',
            'button[class*="model"]',
            'div[class*="model-selector"] button',
        ],
        # Individual model options in the dropdown
        "option_selector": 'div[role="option"], div[role="menuitem"], li[role="option"]',
    },
    "claude": {
        "picker_button": [
            'button[data-testid="model-selector"]',
            'button[class*="model"]',
            'fieldset button[class*="selector"]',
        ],
        "option_selector": 'div[role="option"], div[role="menuitem"], button[role="menuitem"]',
    },
    "gemini": {
        "picker_button": [
            'button[data-test-id="model-selector"]',
            'mat-select[aria-label*="model" i]',
            'button[aria-label*="model" i]',
        ],
        "option_selector": 'mat-option, div[role="option"], div[role="menuitem"]',
    },
}

# Classify follow-up questions: context questions need human input,
# knowledge questions can be routed to other providers.
_CONTEXT_QUESTION_PATTERNS = re.compile(
    r"(which (file|part|section|version|option|approach|language|framework)|"
    r"what (language|framework|version|format|style|tone)|"
    r"your (preference|goal|use case|requirement|deadline|budget)|"
    r"do you (prefer|want|need|mean|have|use)|"
    r"are you (using|looking|trying|working)|"
    r"can you (provide|share|specify|clarify|confirm)|"
    r"what (do you|would you|should I)|"
    r"how (do you want|would you like|should I))",
    re.IGNORECASE,
)


def _color(provider: str, text: str) -> str:
    c = _COLORS.get(provider, "")
    return f"{c}{text}{_RESET}" if c else text


def _open_path(path: Path) -> None:
    """Open a file or directory in the OS default handler (cross-platform)."""
    import subprocess
    try:
        if sys.platform == "darwin":
            subprocess.Popen(["open", str(path)])
        elif sys.platform == "win32":
            os.startfile(str(path))  # type: ignore[attr-defined]
        else:
            subprocess.Popen(["xdg-open", str(path)])
    except Exception:
        pass  # non-critical — user can open manually


def _resolve_file_refs(text: str, cwd: str) -> str:
    """Replace @filename references with file contents.

    Supports:
      @filename.txt       — single file
      @*.md               — glob pattern
      @docs/              — all files in directory
    """
    import glob as globmod

    def _replace(match):
        pattern = match.group(1)
        exact = os.path.join(cwd, pattern)

        # If pattern ends with /, treat as directory — include all files in it
        if pattern.endswith("/") and os.path.isdir(exact):
            candidates = sorted(
                f for f in globmod.glob(os.path.join(exact, "*"))
                if os.path.isfile(f) and not os.path.basename(f).startswith(".")
            )
        elif os.path.isfile(exact):
            candidates = [exact]
        elif os.path.isdir(exact):
            # Bare directory name — include all files
            candidates = sorted(
                f for f in globmod.glob(os.path.join(exact, "*"))
                if os.path.isfile(f) and not os.path.basename(f).startswith(".")
            )
        else:
            candidates = sorted(globmod.glob(os.path.join(cwd, pattern)))
            candidates = [c for c in candidates if os.path.isfile(c)]

        if not candidates:
            return match.group(0)

        parts = []
        for path in candidates[:10]:
            rel = os.path.relpath(path, cwd)
            ext = os.path.splitext(path)[1].lower()
            # Skip files that will be uploaded directly (images, PDFs, etc.)
            if ext in UPLOAD_EXTENSIONS:
                parts.append(f"--- @{rel} --- (will be uploaded directly)")
                continue
            content = _read_file_as_text(path)
            if content is None:
                parts.append(f"--- @{rel} --- (unsupported file type)")
                continue
            if len(content) > MAX_FILE_CHARS:
                content = content[:MAX_FILE_CHARS] + f"\n... (truncated, {len(content)} chars total)"
            parts.append(f"--- @{rel} ---\n{content}\n--- end @{rel} ---")

        return "\n".join(parts)

    result = FILE_REF_PATTERN.sub(_replace, text)
    # If we expanded any @refs, prepend the file context prefix
    if result != text:
        result = _FILE_CONTEXT_PREFIX + result
    return result


# Two sets of keywords — if both sets appear anywhere in the text, auto-include files
_FILE_ACTION_WORDS = re.compile(
    r"\b(review|check|look at|read|analyze|improve|edit|fix|update|refine|critique|evaluate|assess|audit|polish|perfect)\b",
    re.IGNORECASE,
)
_FILE_OBJECT_WORDS = re.compile(
    r"\b(files?|docs?|documents?|scripts?|code|directory|folder|plans?|spreadsheets?|this directory|these)\b",
    re.IGNORECASE,
)
# Direct file extension mentions
_FILE_EXT_PATTERN = re.compile(r"\.(xlsx|docx|md|txt|py|csv|json|pdf|html|pptx)\b", re.IGNORECASE)

TEXT_EXTENSIONS = {
    ".py", ".js", ".ts", ".md", ".txt", ".html", ".css", ".json", ".yaml",
    ".yml", ".toml", ".sh", ".rs", ".go", ".java", ".c", ".cpp", ".h",
    ".rb", ".php", ".sql", ".xml", ".csv", ".ini", ".cfg", ".env",
    ".jsx", ".tsx", ".vue", ".svelte", ".scss", ".less",
}

# Binary formats we can convert to text
BINARY_EXTENSIONS = {".xlsx", ".docx", ".pptx"}
ALL_READABLE = TEXT_EXTENSIONS | BINARY_EXTENSIONS

# Files that should be uploaded directly (not text-converted)
UPLOAD_EXTENSIONS = {
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg", ".bmp", ".ico",
    ".pdf",
    ".zip", ".tar", ".gz",
    ".mp3", ".wav", ".mp4", ".mov",
}

# File input selectors for each provider (hidden <input type="file">)
FILE_INPUT_SELECTORS = {
    "chatgpt": [
        'input[type="file"]',
        'input[data-testid="file-upload"]',
    ],
    "gemini": [
        'input[type="file"]',
    ],
    "claude": [
        'input[type="file"]',
    ],
}


def _read_file_as_text(path: str) -> str | None:
    """Read a file and return its text content. Handles binary formats."""
    ext = os.path.splitext(path)[1].lower()

    if ext in TEXT_EXTENSIONS:
        try:
            content = open(path, "r", errors="replace").read()
            return content
        except Exception:
            return None

    if ext == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"[Sheet: {sheet_name}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    parts.append(" | ".join(cells))
            wb.close()
            return "\n".join(parts)
        except ImportError:
            return f"(Excel file — install openpyxl to read: pip install openpyxl)"
        except Exception as e:
            return f"(Error reading Excel: {e})"

    if ext == ".docx":
        try:
            import docx
            doc = docx.Document(path)
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            return f"(Word file — install python-docx to read: pip install python-docx)"
        except Exception as e:
            return f"(Error reading Word doc: {e})"

    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
                if texts:
                    parts.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(parts)
        except ImportError:
            return f"(PowerPoint file — install python-pptx to read: pip install python-pptx)"
        except Exception as e:
            return f"(Error reading PowerPoint: {e})"

    return None


def _collect_upload_files(text: str, cwd: str) -> list[str]:
    """Find @-referenced files that should be uploaded (images, PDFs, etc.)."""
    upload_paths = []
    for match in FILE_REF_PATTERN.finditer(text):
        pattern = match.group(1)
        exact = os.path.join(cwd, pattern)
        if os.path.isfile(exact):
            ext = os.path.splitext(exact)[1].lower()
            if ext in UPLOAD_EXTENSIONS:
                upload_paths.append(os.path.abspath(exact))
    return upload_paths


def _auto_include_files(text: str, cwd: str) -> str:
    """If user mentions files/docs but didn't use @, auto-include files from cwd."""
    if "@" in text:
        return text  # Already using @ syntax

    # Check if user's intent involves files
    has_action = _FILE_ACTION_WORDS.search(text)
    has_object = _FILE_OBJECT_WORDS.search(text)
    has_ext = _FILE_EXT_PATTERN.search(text)

    if not ((has_action and has_object) or has_ext):
        return text

    # Find readable files in cwd
    files = []
    for f in sorted(os.listdir(cwd)):
        path = os.path.join(cwd, f)
        if os.path.isfile(path) and not f.startswith("."):
            ext = os.path.splitext(f)[1].lower()
            if ext in ALL_READABLE:
                files.append(f)

    if not files:
        return text

    # Auto-include the files with the message
    parts = [_FILE_CONTEXT_PREFIX + text, "\n\nHere are the files from my working directory:\n"]
    for f in files[:10]:
        path = os.path.join(cwd, f)
        content = _read_file_as_text(path)
        if content is None:
            continue
        if len(content) > MAX_FILE_CHARS:
            content = content[:MAX_FILE_CHARS] + f"\n... (truncated, {len(content)} chars total)"
        parts.append(f"\n--- {f} ---\n{content}\n--- end {f} ---")

    return "\n".join(parts)


@dataclass
class BrowserSession:
    """A single browser tab connected to a chat provider."""
    provider: str
    page: object  # playwright Page
    selectors: dict
    messages: list[dict] = field(default_factory=list)
    last_response: str = ""
    call_count: int = 0
    logged_in: bool = False
    # Rate limiting: track last send time and consecutive rate-limit signals
    _last_send_time: float = 0.0
    _rate_limit_hits: int = 0
    # Artifacts extracted from this provider's responses
    artifacts: list[dict] = field(default_factory=list)


def _init_judge_client():
    """Create a free-tier API client for judging/scoring.

    Tries Gemini (free), then Groq (free), then any available provider.
    Returns (client, model, provider_name) or (None, None, None).
    """
    from attestdb.core.providers import PROVIDERS, load_env_file

    # Load .env if present
    env_vars = load_env_file(os.path.join(os.getcwd(), ".env"))

    # Order: free-tier first
    judge_preference = ["gemini", "groq", "deepseek", "openai", "grok", "together"]

    for name in judge_preference:
        provider = PROVIDERS.get(name)
        if not provider:
            continue
        api_key = os.environ.get(provider["env_key"]) or env_vars.get(provider["env_key"])
        if not api_key:
            continue
        try:
            from openai import OpenAI
            client = OpenAI(api_key=api_key, base_url=provider["base_url"])
            return client, provider["default_model"], name
        except Exception:
            continue

    return None, None, None


class BrowserChat:
    """Interactive multi-LLM chat using browser automation.

    Opens ChatGPT, Claude, and Gemini in Chromium tabs using the user's
    saved login sessions. No API keys needed for the primary providers.
    A free-tier API (Gemini/Groq) handles judging and synthesis.
    """

    def __init__(
        self,
        db=None,
        providers: list[str] | None = None,
        cwd: str | None = None,
        profile_dir: str | None = None,
        headless: bool = False,
        models: dict[str, str] | None = None,
        auto_followup: bool = False,
    ):
        self.db = db
        self.cwd = cwd or os.getcwd()
        self.profile_dir = Path(profile_dir) if profile_dir else PROFILE_DIR
        self.headless = headless
        self.requested_providers = providers or ["chatgpt", "claude", "gemini"]
        self.sessions: dict[str, BrowserSession] = {}
        self._context = None
        self._playwright = None
        self._round_num = 0  # Tracks tournament rounds for output dir naming

        # Model selection: {"chatgpt": "o3", "claude": "opus", "gemini": "pro"}
        self._requested_models = models or {}
        self._active_models: dict[str, str] = {}  # Tracks what's actually selected

        # Auto follow-up: route provider questions to other providers before human
        self.auto_followup = auto_followup

        # Output directory for saving responses
        self._output_dir: Path | None = None

        # Judge client (free API for scoring/synthesis)
        self._judge_client = None
        self._judge_model = None
        self._judge_name = None

    def _init_output_dir(self):
        """Create a timestamped output directory for this session."""
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M")
        self._output_dir = Path(self.cwd) / "output" / timestamp
        self._output_dir.mkdir(parents=True, exist_ok=True)
        print(_color("system", f"Saving outputs to: {self._output_dir}/"))

    def _save_response(self, provider: str, content: str, stage: str = "initial"):
        """Save a provider's response to the output directory.

        stage: 'initial', 'round-N', 'consensus', 'final'
        """
        if not self._output_dir or not content:
            return
        stage_dir = self._output_dir / stage
        stage_dir.mkdir(parents=True, exist_ok=True)
        filepath = stage_dir / f"{provider}.md"
        filepath.write_text(content, encoding="utf-8")

    def _save_scores(self, scores_by_judge: dict, stage: str = "round-1"):
        """Save tournament scores to JSON."""
        if not self._output_dir or not scores_by_judge:
            return
        stage_dir = self._output_dir / stage
        stage_dir.mkdir(parents=True, exist_ok=True)
        filepath = stage_dir / "scores.json"
        filepath.write_text(json.dumps(scores_by_judge, indent=2), encoding="utf-8")

    def _save_final(self, winner: str, content: str, metadata: dict | None = None):
        """Save the final consensus/tournament winner."""
        if not self._output_dir or not content:
            return
        final_dir = self._output_dir / "final"
        final_dir.mkdir(parents=True, exist_ok=True)
        (final_dir / "best.md").write_text(
            f"# Best Answer (by {winner})\n\n{content}", encoding="utf-8",
        )
        if metadata:
            (final_dir / "metadata.json").write_text(
                json.dumps(metadata, indent=2), encoding="utf-8",
            )
        print(_color("system", f"  Final output saved to: {final_dir}/best.md"))

    async def _start_browser(self):
        """Launch persistent Chromium context with saved logins.

        Opens all provider tabs in parallel for faster startup.
        Only prompts for login if needed; minimizes the window when all are ready.
        """
        try:
            from playwright.async_api import async_playwright
        except ImportError:
            print("Playwright required. Install with:")
            print("  pip install playwright && playwright install chromium")
            sys.exit(1)

        self.profile_dir.mkdir(parents=True, exist_ok=True)

        self._playwright = await async_playwright().start()
        self._context = await self._playwright.chromium.launch_persistent_context(
            str(self.profile_dir),
            headless=self.headless,
            accept_downloads=True,
            viewport={"width": 1280, "height": 900},
            args=["--disable-blink-features=AutomationControlled"],
        )

        # Close the default blank tab immediately
        for p in self._context.pages:
            if p.url in ("about:blank", "chrome://newtab/"):
                await p.close()

        # Open ALL tabs in parallel
        valid_providers = [p for p in self.requested_providers if p in PROVIDER_SELECTORS]
        print(_color("system", f"Opening {len(valid_providers)} providers in parallel..."))

        async def _open_tab(provider: str):
            selectors = PROVIDER_SELECTORS[provider]
            page = await self._context.new_page()
            try:
                await page.goto(selectors["url"], wait_until="domcontentloaded", timeout=30000)
                # Wait for redirects / consent screens to settle
                try:
                    await page.wait_for_load_state("networkidle", timeout=10000)
                except Exception:
                    await asyncio.sleep(2)
            except Exception as exc:
                print(_color("system", f"  Failed to load {provider}: {exc}"))
                return None
            return BrowserSession(
                provider=provider, page=page, selectors=selectors,
            )

        results = await asyncio.gather(
            *[_open_tab(p) for p in valid_providers],
            return_exceptions=True,
        )

        # Collect successful sessions, track which need login
        needs_login = []
        for provider, result in zip(valid_providers, results):
            if isinstance(result, Exception) or result is None:
                print(_color("system", f"  {provider}: failed to open"))
                continue
            session = result
            logged_in = await self._check_login(session)
            if logged_in:
                session.logged_in = True
                self.sessions[provider] = session
                print(_color("system", f"  {provider} ✓"))
            else:
                needs_login.append(session)

        # Handle logins one at a time — bring each tab to front
        if needs_login:
            print()
            print(_color("system", f"{len(needs_login)} provider(s) need login:"))
            for session in needs_login:
                provider = session.provider
                try:
                    await session.page.bring_to_front()
                except Exception:
                    pass
                print(f"\n  {_BOLD}Log into {provider} in the browser window, then press Enter here.{_RESET}")
                await asyncio.get_event_loop().run_in_executor(None, input, "  Press Enter when logged in... ")
                await asyncio.sleep(2)

                # Re-navigate to ensure we're on the chat page
                try:
                    await session.page.goto(
                        session.selectors["url"],
                        wait_until="domcontentloaded", timeout=15000,
                    )
                    await asyncio.sleep(1)
                except Exception:
                    pass

                session.logged_in = True
                self.sessions[provider] = session
                print(_color("system", f"  {provider} ✓"))

        # Register download handlers
        for session in self.sessions.values():
            await self._setup_download_handler(session)

        # Select requested models (if any)
        if self._requested_models:
            print(_color("system", "Selecting models..."))
            for provider, model_key in self._requested_models.items():
                if provider in self.sessions:
                    await self._select_model(self.sessions[provider], model_key)

        # Minimize the browser window — work happens in the terminal.
        # Tabs will be brought to front only when follow-up input is needed.
        if self.sessions and not needs_login:
            try:
                first_page = next(iter(self.sessions.values())).page
                await first_page.evaluate("() => window.blur()")
            except Exception:
                pass

    async def _check_login(self, session: BrowserSession) -> bool:
        """Heuristic check if the user is logged into this provider."""
        page = session.page
        provider = session.provider

        try:
            if provider == "chatgpt":
                # If we see the prompt input, we're logged in
                el = await page.query_selector(session.selectors["input"])
                return el is not None
            elif provider == "claude":
                el = await page.query_selector(session.selectors["input"])
                return el is not None
            elif provider == "gemini":
                el = await page.query_selector(session.selectors["input"])
                return el is not None
        except Exception:
            pass
        return False

    async def _select_model(self, session: BrowserSession, model_key: str) -> bool:
        """Select a specific model/tier in a provider's web UI.

        model_key: short name like "opus", "o3", "pro", "deep-research"
        Returns True if selection succeeded (or was already active).
        """
        provider = session.provider
        page = session.page

        models = PROVIDER_MODELS.get(provider, {})
        if model_key not in models:
            available = ", ".join(models.keys()) if models else "none"
            print(_color("system", f"  Unknown model '{model_key}' for {provider}. Available: {available}"))
            return False

        display_text = models[model_key]
        selectors = MODEL_SELECTORS.get(provider, {})
        picker_selectors = selectors.get("picker_button", [])
        option_selector = selectors.get("option_selector", 'div[role="option"]')

        # Step 1: Click the model picker to open dropdown
        picker = None
        for sel in picker_selectors:
            try:
                picker = await page.query_selector(sel)
                if picker:
                    break
            except Exception:
                continue

        if not picker:
            # Try a broader search — look for any button containing current model text
            try:
                picker = await page.evaluate("""() => {
                    const buttons = document.querySelectorAll('button');
                    for (const btn of buttons) {
                        const text = btn.innerText || btn.textContent || '';
                        if (/GPT|Claude|Sonnet|Opus|Haiku|Gemini|Flash|Pro|model/i.test(text)) {
                            return true;  // found it
                        }
                    }
                    return false;
                }""")
                if picker:
                    # Click it via JS
                    await page.evaluate("""() => {
                        const buttons = document.querySelectorAll('button');
                        for (const btn of buttons) {
                            const text = btn.innerText || btn.textContent || '';
                            if (/GPT|Claude|Sonnet|Opus|Haiku|Gemini|Flash|Pro|model/i.test(text)) {
                                btn.click();
                                return true;
                            }
                        }
                        return false;
                    }""")
            except Exception:
                pass

        if not picker:
            print(_color("system", f"  Could not find model picker for {provider}"))
            return False

        if isinstance(picker, bool):
            # JS-based click already happened
            pass
        else:
            try:
                await picker.click()
            except Exception as exc:
                print(_color("system", f"  Could not click model picker for {provider}: {exc}"))
                return False

        await asyncio.sleep(1.0)  # Wait for dropdown to open

        # Step 2: Find and click the target model option
        try:
            clicked = await page.evaluate(f"""(targetText) => {{
                const options = document.querySelectorAll('{option_selector}');
                for (const opt of options) {{
                    const text = opt.innerText || opt.textContent || '';
                    if (text.includes(targetText)) {{
                        opt.click();
                        return text.trim();
                    }}
                }}
                // Broader fallback: any clickable element containing the text
                const all = document.querySelectorAll('div, span, button, li');
                for (const el of all) {{
                    const text = el.innerText || el.textContent || '';
                    if (text.includes(targetText) && el.offsetParent !== null) {{
                        el.click();
                        return text.trim();
                    }}
                }}
                return null;
            }}""", display_text)

            if clicked:
                self._active_models[provider] = model_key
                print(_color("system", f"  {provider}: selected {clicked}"))
                await asyncio.sleep(1.0)  # Wait for model switch to take effect
                return True
            else:
                print(_color("system", f"  Could not find '{display_text}' option in {provider} dropdown"))
                # Close the dropdown by pressing Escape
                await page.keyboard.press("Escape")
                return False
        except Exception as exc:
            print(_color("system", f"  Model selection failed for {provider}: {exc}"))
            await page.keyboard.press("Escape")
            return False

    def _is_context_question(self, response: str) -> bool:
        """Check if a follow-up question is about the user's specific context
        (needs human input) vs a knowledge question (other LLMs can answer)."""
        # Extract the question portion (last part of response)
        tail = response[-1000:] if len(response) > 1000 else response
        return bool(_CONTEXT_QUESTION_PATTERNS.search(tail))

    async def _route_to_other_providers(
        self, asking_provider: str, question_text: str,
    ) -> str | None:
        """Send a provider's follow-up question to the other providers.

        Returns a synthesized answer, or None if no useful responses.
        """
        other_sessions = {
            name: session for name, session in self.sessions.items()
            if name != asking_provider and session.logged_in
        }
        if not other_sessions:
            return None

        # Frame the question clearly
        routed_msg = (
            f"Another AI assistant ({asking_provider}) asked this question while "
            f"working on the same task. Please answer briefly:\n\n"
            f"{question_text}"
        )

        print(_color("system", f"  Routing {asking_provider}'s question to: {', '.join(other_sessions.keys())}"))

        # Send to all other providers in parallel
        task_to_provider = {}
        pending = set()
        for name, session in other_sessions.items():
            task = asyncio.create_task(self._send_message(session, routed_msg))
            task_to_provider[task] = name
            pending.add(task)

        answers = {}
        while pending:
            done, pending = await asyncio.wait(pending, return_when=asyncio.FIRST_COMPLETED)
            for done_task in done:
                name = task_to_provider[done_task]
                try:
                    resp = done_task.result()
                    if resp and len(resp.strip()) > 10:
                        answers[name] = resp
                except Exception:
                    pass

        if not answers:
            return None

        # If only one answer, use it directly
        if len(answers) == 1:
            provider, answer = next(iter(answers.items()))
            print(_color("system", f"  {provider} answered {asking_provider}'s question"))
            return f"[Answer from {provider}]: {answer}"

        # Multiple answers — synthesize
        parts = [f"Here are answers from other AI assistants to your question:"]
        for name, answer in answers.items():
            # Truncate individual answers to keep context manageable
            truncated = answer[:1500] if len(answer) > 1500 else answer
            parts.append(f"\n[{name}]: {truncated}")

        parts.append(
            "\nBased on these answers, please continue with your original task."
        )
        return "\n".join(parts)

    async def _check_page_health(self, session: BrowserSession) -> bool:
        """Check if a provider's page is still alive and on the right site.

        If the page has crashed or navigated away (session expired, etc.),
        attempt to recover by navigating back to the provider's URL.
        """
        page = session.page
        provider = session.provider

        try:
            url = page.url
            if not url or url == "about:blank" or url.startswith("chrome-error"):
                raise Exception("Page crashed or blank")
            # Check the page is on the right domain
            expected_domain = PROVIDER_SELECTORS[provider]["url"].split("//")[1].split("/")[0]
            if expected_domain not in url:
                raise Exception(f"Wrong domain: {url}")
            return True
        except Exception as exc:
            logger.warning("%s page unhealthy: %s — attempting recovery", provider, exc)
            print(_color("system", f"  {provider}: page issue detected, recovering..."))
            try:
                await page.goto(
                    session.selectors["url"],
                    wait_until="domcontentloaded", timeout=15000,
                )
                await asyncio.sleep(3)
                # Verify input is available
                el = await self._find_element(
                    page, session.selectors["input"],
                    session.selectors.get("input_fallbacks"), timeout=5000,
                )
                if el:
                    print(_color("system", f"  {provider}: recovered"))
                    return True
                else:
                    print(_color("system", f"  {provider}: recovered page but can't find input — may need re-login"))
                    return False
            except Exception as exc2:
                print(_color("system", f"  {provider}: recovery failed: {exc2}"))
                return False

    async def _find_element(self, page, primary: str, fallbacks: list[str] | None = None, timeout: int = 10000) -> object | None:
        """Find an element using primary selector with fallbacks."""
        try:
            el = await page.wait_for_selector(primary, timeout=timeout)
            if el:
                return el
        except Exception:
            pass

        for sel in (fallbacks or []):
            try:
                el = await page.query_selector(sel)
                if el:
                    return el
            except Exception:
                continue
        return None

    async def _insert_text(self, page, element, text: str, provider: str):
        """Insert text into a chat input as a single atomic paste.

        All providers use the same approach: focus the element, write to the
        system clipboard via JS, then trigger a real paste event via Ctrl/Cmd+V.
        This is the most reliable method because it's exactly what a human does
        and every editor framework handles paste as a single atomic operation.
        """
        await element.click()
        await asyncio.sleep(0.3)

        # Clear any existing text
        await page.keyboard.press("Meta+a")
        await asyncio.sleep(0.1)
        await page.keyboard.press("Backspace")
        await asyncio.sleep(0.1)

        # Strategy 1: Clipboard API paste (works for all providers)
        try:
            # Write text to clipboard via the Clipboard API
            await page.evaluate("""(text) => {
                // Use a synthetic paste event with DataTransfer
                // This works even without clipboard permissions
                const el = document.activeElement;
                if (el) {
                    const dt = new DataTransfer();
                    dt.setData('text/plain', text);
                    const event = new ClipboardEvent('paste', {
                        clipboardData: dt,
                        bubbles: true,
                        cancelable: true,
                    });
                    el.dispatchEvent(event);
                }
            }""", text)
            await asyncio.sleep(0.5)

            # Check if the text actually landed
            content = await self._get_input_content(page, provider)
            if content and len(content) > min(20, len(text) // 2):
                return
        except Exception:
            pass

        # Strategy 2: Clipboard write + real Cmd+V keypress
        try:
            # Grant clipboard permissions and write
            await page.context.grant_permissions(["clipboard-read", "clipboard-write"])
        except Exception:
            pass

        try:
            await page.evaluate("(text) => navigator.clipboard.writeText(text)", text)
            await asyncio.sleep(0.2)
            await page.keyboard.press("Meta+v")
            await asyncio.sleep(0.5)

            content = await self._get_input_content(page, provider)
            if content and len(content) > min(20, len(text) // 2):
                return
        except Exception:
            pass

        # Strategy 3: Provider-specific innerHTML (fast but less reliable)
        if provider == "chatgpt":
            try:
                paragraphs = text.split("\n")
                html = "".join(f"<p>{p}</p>" if p.strip() else "<p><br></p>" for p in paragraphs)
                await page.evaluate("""([selector, html]) => {
                    const el = document.querySelector(selector);
                    if (el) {
                        el.innerHTML = html;
                        // Trigger React's synthetic event system
                        const inputEvent = new InputEvent('input', { bubbles: true, inputType: 'insertFromPaste' });
                        el.dispatchEvent(inputEvent);
                    }
                }""", ['div[id="prompt-textarea"]', html])
                await asyncio.sleep(0.5)
                return
            except Exception:
                pass
        elif provider == "gemini":
            try:
                await page.evaluate("""([selector, text]) => {
                    const el = document.querySelector(selector);
                    if (el) {
                        el.innerHTML = '<p>' + text.replace(/\\n/g, '</p><p>') + '</p>';
                        const inputEvent = new InputEvent('input', { bubbles: true, inputType: 'insertFromPaste' });
                        el.dispatchEvent(inputEvent);
                    }
                }""", ['div.ql-editor', text])
                await asyncio.sleep(0.5)
                return
            except Exception:
                pass

        # Strategy 4: keyboard.insert_text (Playwright's built-in, inserts all at once)
        try:
            await element.click()
            await asyncio.sleep(0.1)
            await page.keyboard.insert_text(text)
            await asyncio.sleep(0.3)
            return
        except Exception:
            pass

        # Strategy 5: Last resort — type it (slow but guaranteed)
        logger.warning("All paste strategies failed for %s, falling back to typing", provider)
        truncated = text[:5000] if len(text) > 5000 else text
        await page.keyboard.type(truncated, delay=1)

    async def _get_input_content(self, page, provider: str) -> str | None:
        """Check what text is currently in the input element."""
        try:
            selectors = {
                "chatgpt": 'div[id="prompt-textarea"]',
                "gemini": 'div.ql-editor',
                "claude": 'div[contenteditable="true"]',
            }
            sel = selectors.get(provider, 'div[contenteditable="true"]')
            content = await page.evaluate(f"""() => {{
                const el = document.querySelector('{sel}');
                return el ? el.innerText : null;
            }}""")
            return content
        except Exception:
            return None

    async def _upload_files(self, session: BrowserSession, file_paths: list[str]) -> int:
        """Upload files to a provider's chat UI using the hidden file input.

        Returns number of files successfully uploaded.
        """
        if not file_paths:
            return 0

        page = session.page
        provider = session.provider
        selectors = FILE_INPUT_SELECTORS.get(provider, ['input[type="file"]'])

        uploaded = 0
        for sel in selectors:
            try:
                file_input = await page.query_selector(sel)
                if file_input:
                    await file_input.set_input_files(file_paths)
                    uploaded = len(file_paths)
                    await asyncio.sleep(2)  # Wait for upload processing
                    break
            except Exception as exc:
                logger.debug("Upload selector %s failed for %s: %s", sel, provider, exc)
                continue

        if uploaded == 0:
            # Fallback: trigger file chooser dialog via the attach button
            try:
                attach_selectors = {
                    "chatgpt": ['button[aria-label="Attach files"]', 'button[aria-label="Upload file"]'],
                    "gemini": ['button[aria-label="Upload file"]', 'button[aria-label="Add files"]', 'button[mattooltip="Upload file"]'],
                    "claude": ['button[aria-label="Attach files"]', 'button[aria-label="Upload files"]'],
                }
                for btn_sel in attach_selectors.get(provider, []):
                    try:
                        async with page.expect_file_chooser(timeout=3000) as fc_info:
                            btn = await page.query_selector(btn_sel)
                            if btn:
                                await btn.click()
                        file_chooser = await fc_info.value
                        await file_chooser.set_files(file_paths)
                        uploaded = len(file_paths)
                        await asyncio.sleep(2)
                        break
                    except Exception:
                        continue
            except Exception as exc:
                logger.debug("File chooser fallback failed for %s: %s", provider, exc)

        if uploaded:
            names = [os.path.basename(f) for f in file_paths]
            print(_color("system", f"  Uploaded {uploaded} file(s) to {provider}: {', '.join(names)}"))
        else:
            print(_color("system", f"  Could not upload files to {provider} (no file input found)"))

        return uploaded

    async def _setup_download_handler(self, session: BrowserSession):
        """Register a download handler for a provider page.

        Downloads are saved to the output directory under downloads/.
        """
        page = session.page
        provider = session.provider

        async def on_download(download):
            if not self._output_dir:
                return
            dl_dir = self._output_dir / "downloads" / provider
            dl_dir.mkdir(parents=True, exist_ok=True)
            filename = download.suggested_filename
            save_path = dl_dir / filename
            try:
                await download.save_as(str(save_path))
                print(_color("system", f"  Downloaded from {provider}: {filename} → {save_path}"))
            except Exception as exc:
                print(_color("system", f"  Download failed from {provider}: {exc}"))

        page.on("download", on_download)

    async def _count_responses(self, session: BrowserSession) -> int:
        """Count current number of assistant response elements."""
        page = session.page
        sel = session.selectors
        all_selectors = [sel["response"]] + sel.get("response_fallbacks", [])
        for s in all_selectors:
            try:
                els = await page.query_selector_all(s)
                if els:
                    return len(els)
            except Exception:
                continue
        return 0

    # Rate-limit detection patterns (provider shows a "slow down" or "try again" message)
    _RATE_LIMIT_PATTERNS = re.compile(
        r"(rate limit|too many (requests|messages)|slow down|try again (later|in)|"
        r"wait a (moment|few|bit)|usage cap|limit reached|you.ve sent too many)",
        re.IGNORECASE,
    )

    async def _rate_limit_delay(self, session: BrowserSession):
        """Enforce per-provider cooldown between sends.

        Base delay: 2s between messages. If rate-limit signals detected,
        exponential backoff up to 60s.
        """
        now = time.monotonic()
        base_delay = 2.0
        if session._rate_limit_hits > 0:
            # Exponential backoff: 5s, 10s, 20s, 40s, 60s cap
            delay = min(5.0 * (2 ** (session._rate_limit_hits - 1)), 60.0)
            print(_color("system", f"  {session.provider}: rate-limit backoff {delay:.0f}s"))
        else:
            delay = base_delay

        elapsed = now - session._last_send_time
        if elapsed < delay:
            wait = delay - elapsed
            await asyncio.sleep(wait)

        session._last_send_time = time.monotonic()

    def _check_rate_limit_response(self, session: BrowserSession, response: str | None):
        """Check if a response indicates rate limiting and update the session."""
        if not response:
            return
        if self._RATE_LIMIT_PATTERNS.search(response) and len(response) < 500:
            session._rate_limit_hits = min(session._rate_limit_hits + 1, 5)
            logger.warning("%s: rate limit detected (hit #%d)", session.provider, session._rate_limit_hits)
        else:
            # Successful response — decay rate limit counter
            if session._rate_limit_hits > 0:
                session._rate_limit_hits = max(0, session._rate_limit_hits - 1)

    async def _send_message(self, session: BrowserSession, text: str, stream_display: bool = False) -> str | None:
        """Type and send a message, wait for response. Returns response text.

        If stream_display=True, prints the response incrementally as it streams.
        Only use for single-provider sends to avoid interleaving output.
        """
        page = session.page
        sel = session.selectors

        # Rate limit delay before sending
        await self._rate_limit_delay(session)

        try:
            # Snapshot response count before sending
            pre_count = await self._count_responses(session)

            # Find input element with fallbacks
            input_el = await self._find_element(
                page, sel["input"], sel.get("input_fallbacks"), timeout=10000,
            )
            if not input_el:
                logger.error("Cannot find input for %s", session.provider)
                print(_color("system", f"  {session.provider}: cannot find input box"))
                return None

            # Insert text (provider-specific strategy)
            await self._insert_text(page, input_el, text, session.provider)
            await asyncio.sleep(2.0)  # Let the editor framework fully process

            # Verify text landed before sending
            content = await self._get_input_content(page, session.provider)
            if not content or len(content.strip()) < min(10, len(text) // 4):
                logger.warning("%s: text may not have landed (got %d chars), retrying",
                               session.provider, len(content or ""))
                # Retry with a different strategy — click and insert_text
                await input_el.click()
                await asyncio.sleep(0.2)
                await page.keyboard.press("Meta+a")
                await asyncio.sleep(0.1)
                await page.keyboard.insert_text(text)
                await asyncio.sleep(1.0)

            # Click send button (with fallbacks), or press Enter
            send_btn = await self._find_element(
                page, sel["send"], sel.get("send_fallbacks"), timeout=5000,
            )
            sent = False
            if send_btn:
                try:
                    await send_btn.click()
                    sent = True
                except Exception:
                    pass

            if not sent:
                # Try Enter key, then Ctrl+Enter (some UIs use this)
                await page.keyboard.press("Enter")

            # Wait for streaming to start — longer for big messages
            start_wait = 5 if len(text) > 2000 else 3
            await asyncio.sleep(start_wait)

            # Verify the message was sent: check if a new user message appeared
            # or if streaming started. If neither, try re-sending via Enter.
            streaming = await self._is_streaming(session)
            new_count = await self._count_responses(session)
            if not streaming and new_count <= pre_count:
                logger.info("%s: message may not have sent, retrying Enter", session.provider)
                # Focus the input again and try Enter
                try:
                    input_el2 = await self._find_element(
                        page, sel["input"], sel.get("input_fallbacks"), timeout=3000,
                    )
                    if input_el2:
                        await input_el2.click()
                        await asyncio.sleep(0.3)
                except Exception:
                    pass
                await page.keyboard.press("Enter")
                await asyncio.sleep(3)

            # Wait for streaming to finish
            response = await self._wait_for_response(session, pre_count=pre_count, stream_display=stream_display)

            # Check for rate-limiting in the response
            self._check_rate_limit_response(session, response)

            if response:
                session.last_response = response
                session.call_count += 1
                session.messages.append({"role": "user", "content": text})
                session.messages.append({"role": "assistant", "content": response})

                # Extract artifacts (code blocks, canvas, etc.)
                try:
                    artifacts = await self._extract_artifacts(session)
                    if artifacts:
                        # Store on session and save to disk
                        session.artifacts.extend(artifacts)
                        self._save_artifacts(session.provider, artifacts)
                except Exception:
                    pass  # Artifact extraction is best-effort

            return response

        except Exception as exc:
            print(_color("system", f"  Failed to send to {session.provider}: {exc}"))
            logger.error("Failed to send to %s: %s", session.provider, exc)
            return None

    async def _is_streaming(self, session: BrowserSession) -> bool:
        """Check if the provider is still generating a response."""
        page = session.page
        sel = session.selectors
        all_indicators = [sel.get("streaming_indicator", "nonexistent")] + sel.get("streaming_fallbacks", [])
        for indicator in all_indicators:
            try:
                el = await page.query_selector(indicator)
                if el:
                    return True
            except Exception:
                continue
        return False

    async def _extract_last_response(self, session: BrowserSession, pre_count: int = 0) -> str | None:
        """Extract the last assistant response using primary + fallback selectors.

        Uses pre_count to only consider NEW response elements — prevents picking
        up stale responses from earlier rounds during tournaments.
        """
        page = session.page
        sel = session.selectors
        all_selectors = [sel["response"]] + sel.get("response_fallbacks", [])

        for selector in all_selectors:
            try:
                response_els = await page.query_selector_all(selector)
                if not response_els:
                    continue

                # Only look at elements added AFTER the pre_count snapshot.
                # If pre_count > 0 and we have more elements now, skip old ones.
                if pre_count > 0 and len(response_els) > pre_count:
                    candidates = response_els[pre_count:]
                else:
                    candidates = [response_els[-1]]

                # Get the last NEW element
                for el in reversed(candidates):
                    text = await el.inner_text()
                    if text and text.strip() and len(text.strip()) > 5:
                        return text.strip()
            except Exception:
                continue

        return None

    async def _extract_artifacts(self, session: BrowserSession) -> list[dict]:
        """Extract code blocks and artifacts from the provider's latest response.

        Returns a list of dicts with keys: type, content, language (optional).
        """
        page = session.page
        sel = session.selectors
        artifact_selectors = sel.get("artifact_selectors", [])
        artifacts = []

        for selector in artifact_selectors:
            try:
                elements = await page.query_selector_all(selector)
                for el in elements:
                    text = await el.inner_text()
                    if not text or len(text.strip()) < 10:
                        continue

                    # Try to detect the language from class attributes
                    lang = ""
                    try:
                        class_attr = await el.get_attribute("class") or ""
                        # Common patterns: "language-python", "hljs python", "code-python"
                        for cls in class_attr.split():
                            if cls.startswith("language-"):
                                lang = cls[9:]
                                break
                            elif cls in ("python", "javascript", "typescript", "rust",
                                         "go", "java", "cpp", "c", "html", "css", "sql",
                                         "bash", "shell", "json", "yaml", "xml", "markdown"):
                                lang = cls
                                break
                    except Exception:
                        pass

                    artifact = {
                        "type": "code" if selector.endswith("code") else "artifact",
                        "content": text.strip(),
                        "language": lang,
                        "provider": session.provider,
                    }

                    # Deduplicate — don't add if we already have this exact content
                    if not any(a["content"] == artifact["content"] for a in artifacts):
                        artifacts.append(artifact)

            except Exception:
                continue

        return artifacts

    def _save_artifacts(self, provider: str, artifacts: list[dict], stage: str = "initial"):
        """Save extracted artifacts to the output directory."""
        if not self._output_dir or not artifacts:
            return
        art_dir = self._output_dir / "artifacts" / provider
        art_dir.mkdir(parents=True, exist_ok=True)
        for i, art in enumerate(artifacts):
            ext = art.get("language") or "txt"
            # Map common language names to file extensions
            ext_map = {
                "python": "py", "javascript": "js", "typescript": "ts",
                "rust": "rs", "bash": "sh", "shell": "sh", "markdown": "md",
                "cpp": "cpp", "c": "c", "yaml": "yml",
            }
            ext = ext_map.get(ext, ext)
            filename = f"{stage}_{i+1}.{ext}"
            (art_dir / filename).write_text(art["content"], encoding="utf-8")
        print(_color("system", f"  {provider}: saved {len(artifacts)} artifact(s) to artifacts/{provider}/"))

    def _is_deep_research(self, provider: str) -> bool:
        """Check if a provider is in a deep-research / extended mode."""
        model = self._active_models.get(provider, "")
        return model in ("deep-research", "o1", "o3")

    async def _wait_for_response(self, session: BrowserSession, timeout: float = 180.0, pre_count: int = 0, stream_display: bool = True) -> str | None:
        """Poll until the assistant finishes responding.

        Strategy: always try to extract text every cycle, whether streaming
        or not. Text must stabilize (identical for 2 consecutive checks)
        AND streaming must have stopped.

        If stream_display=True, prints incremental text as it streams in
        so the user can read the response in real-time.

        Deep Research / reasoning models get a 15-minute timeout and periodic
        progress updates in the terminal.
        """
        # Deep Research and reasoning models can take much longer
        if self._is_deep_research(session.provider):
            timeout = max(timeout, 900.0)  # 15 minutes
            logger.info("%s is in deep research mode — timeout set to %ds", session.provider, timeout)
        start = time.monotonic()
        interval = 2.0
        last_text = ""
        stable_count = 0
        found_any_text = False
        last_progress_print = 0.0  # Track when we last printed a progress update
        displayed_len = 0  # How much of the response we've already printed
        stream_started = False  # Have we printed the streaming header?

        while (time.monotonic() - start) < timeout:
            streaming = await self._is_streaming(session)

            # Always try to extract, even while streaming
            text = await self._extract_last_response(session, pre_count)

            if text and len(text) > 5:
                found_any_text = True

                # Stream display: print new characters as they arrive
                if stream_display and streaming and len(text) > displayed_len:
                    if not stream_started:
                        stream_started = True
                        print(_color(session.provider, f"─── {session.provider} (streaming) ───"))
                    new_text = text[displayed_len:]
                    print(new_text, end="", flush=True)
                    displayed_len = len(text)

                if not streaming:
                    # Only check stability when NOT streaming
                    if text == last_text:
                        stable_count += 1
                        if stable_count >= 2:
                            # Print any remaining text that wasn't streamed
                            if stream_display and stream_started:
                                remaining = text[displayed_len:]
                                if remaining:
                                    print(remaining, end="", flush=True)
                                print()  # Newline after streaming
                            return text
                    else:
                        last_text = text
                        stable_count = 1
                else:
                    # Still streaming — track text but reset stability
                    last_text = text
                    stable_count = 0
            elif not found_any_text and (time.monotonic() - start) > 30:
                # 30s with zero text — try a JS-based extraction as last resort
                try:
                    js_text = await session.page.evaluate("""() => {
                        // Try multiple strategies to find response text
                        const selectors = [
                            'div[data-message-author-role="assistant"]',
                            'div.markdown.prose',
                            'div.agent-turn',
                            'article[data-testid] div.markdown',
                            'div.font-claude-message',
                            'message-content',
                        ];
                        for (const sel of selectors) {
                            const els = document.querySelectorAll(sel);
                            if (els.length > 0) {
                                const last = els[els.length - 1];
                                const text = last.innerText;
                                if (text && text.trim().length > 10) return text.trim();
                            }
                        }
                        return null;
                    }""")
                    if js_text:
                        logger.info("JS fallback found text for %s (%d chars)", session.provider, len(js_text))
                        found_any_text = True
                        last_text = js_text
                        stable_count = 1
                except Exception:
                    pass

            # Print progress for deep research / long-running responses
            elapsed = time.monotonic() - start
            if self._is_deep_research(session.provider) and elapsed - last_progress_print > 30:
                last_progress_print = elapsed
                mins = int(elapsed // 60)
                secs = int(elapsed % 60)
                chars = len(last_text) if last_text else 0
                status = "streaming" if streaming else ("waiting" if not found_any_text else "stabilizing")
                print(_color("system", f"  {session.provider}: {status} ({mins}m{secs:02d}s, {chars} chars)"), end="\r")

            await asyncio.sleep(interval)
            interval = min(interval * 1.1, 4.0)

        # Return whatever we have even on timeout
        if last_text:
            logger.warning("Timeout for %s but returning partial text (%d chars)", session.provider, len(last_text))
            return last_text

        logger.warning("Timeout waiting for %s — no text found", session.provider)
        return None

    def _is_stalled(self, response: str) -> bool:
        """Check if a response is a stall (asking for files instead of reviewing)."""
        if not response:
            return False
        if len(response) > _STALL_MAX_LEN:
            return False  # Long responses are probably real content
        return bool(_STALL_PATTERNS.search(response))

    def _needs_followup(self, response: str) -> bool:
        """Check if a response is asking the user follow-up questions."""
        if not response:
            return False
        # Count question marks in the last portion of the response
        tail = response[-500:] if len(response) > 500 else response
        n_questions = tail.count("?")
        if n_questions >= _MIN_QUESTIONS:
            return True
        return bool(_FOLLOWUP_PATTERNS.search(response))

    async def _handle_followup(self, provider: str, session: BrowserSession, response: str, _depth: int = 0) -> str | None:
        """Handle a provider asking follow-up questions.

        Modes:
        - auto_followup: Knowledge questions go to other providers first.
          Context questions still go to the human.
        - Manual (default): Always ask the human, with option to type 'ask'
          to route the question to other providers.

        User options:
        - Type answer    → sends it to the provider
        - Type 'done'    → you already answered in the browser
        - Type 'ask'     → route the question to other providers
        - Press Enter    → tells provider to proceed with best judgment
        """
        # Auto-followup mode: try other providers first for knowledge questions
        if self.auto_followup and _depth == 0 and not self._is_context_question(response):
            # Extract the question portion for routing
            tail = response[-2000:] if len(response) > 2000 else response
            print(_color("system", f"  {provider} has a question — routing to other providers (auto-followup)..."))
            routed_answer = await self._route_to_other_providers(provider, tail)
            if routed_answer:
                new_response = await self._send_message(session, routed_answer)
                if new_response:
                    if _depth < 2 and self._needs_followup(new_response) and not self._is_stalled(new_response):
                        # Still asking questions — fall through to human
                        print(_color("system", f"  {provider} still has questions after auto-followup"))
                    else:
                        return new_response
            else:
                print(_color("system", f"  No other providers available — asking you"))

        # Bring the tab to front immediately so user sees the questions
        try:
            await session.page.bring_to_front()
        except Exception:
            pass

        # Show the questions in terminal too
        print()
        print(_color("system", f"╔══ {provider} has questions (tab activated) ══╗"))
        # Show a condensed version in terminal — user can read full context in browser
        tail = response[-1500:] if len(response) > 1500 else response
        print(_color(provider, tail))
        print()
        print(_color("system", f"  Type answer here  → sends it to {provider}"))
        print(_color("system", f"  Type 'done'       → if you already answered in the browser"))
        print(_color("system", f"  Type 'ask'        → route question to other providers"))
        print(_color("system", f"  Press Enter       → tells {provider} to proceed with best judgment"))
        print()

        user_answer = await asyncio.get_event_loop().run_in_executor(
            None, lambda: input(f"{_BOLD}  reply to {provider}>{_RESET} ").strip()
        )

        if user_answer.lower() == "done":
            # User answered directly in browser — wait for the new response
            print(_color("system", f"  Waiting for {provider} to respond..."))
            await asyncio.sleep(3)
            new_response = await self._wait_for_response(session, timeout=120.0)
            if new_response and new_response != response:
                session.last_response = new_response
                return new_response
            # If no new response, try extracting
            new_response = await self._extract_last_response(session)
            if new_response and new_response != response:
                session.last_response = new_response
                return new_response
            return session.last_response

        if user_answer.lower() == "ask":
            # Route the question to other providers
            tail = response[-2000:] if len(response) > 2000 else response
            routed_answer = await self._route_to_other_providers(provider, tail)
            if routed_answer:
                new_response = await self._send_message(session, routed_answer)
                if new_response:
                    if _depth < 2 and self._needs_followup(new_response) and not self._is_stalled(new_response):
                        print(_color("system", f"  {provider} has more questions..."))
                        return await self._handle_followup(provider, session, new_response, _depth=_depth + 1)
                    return new_response
            else:
                print(_color("system", f"  No other providers available"))

        # Send the answer (or default "proceed" message)
        if not user_answer or user_answer.lower() == "ask":
            user_answer = (
                "Proceed with your best judgment. Don't ask more questions — "
                "just give me your complete, best answer based on what you have."
            )

        new_response = await self._send_message(session, user_answer)
        if new_response:
            # Check if STILL asking questions (cap recursion at 2 levels)
            if _depth < 2 and self._needs_followup(new_response) and not self._is_stalled(new_response):
                print(_color("system", f"  {provider} has more questions..."))
                return await self._handle_followup(provider, session, new_response, _depth=_depth + 1)
        return new_response

    async def _send_to_all(self, message: str, files_included: bool = False, upload_files: list[str] | None = None):
        """Send message to all providers in parallel, print responses.

        If files_included=True and a provider stalls (asks for uploads),
        auto-sends a nudge to get it back on track.
        If upload_files is provided, uploads those files before sending the message.
        """
        print()

        # Check page health and attempt recovery for any broken tabs
        active_sessions = {}
        for provider, session in self.sessions.items():
            healthy = await self._check_page_health(session)
            if healthy:
                active_sessions[provider] = session
            else:
                print(_color("system", f"  {provider}: unavailable, skipping"))

        # Upload files to all providers first (if any)
        if upload_files:
            upload_tasks = {
                provider: asyncio.create_task(self._upload_files(session, upload_files))
                for provider, session in active_sessions.items()
            }
            for provider, task in upload_tasks.items():
                try:
                    await task
                except Exception as exc:
                    print(_color("system", f"  {provider}: upload error: {exc}"))

        # Send to all active providers concurrently
        task_to_provider = {}
        pending = set()
        for provider, session in active_sessions.items():
            task = asyncio.create_task(self._send_message(session, message))
            task_to_provider[task] = provider
            pending.add(task)

        # Process responses as they arrive (fastest first) — follow-ups
        # get handled immediately instead of waiting for all providers
        stalled_providers = []
        while pending:
            done, pending = await asyncio.wait(pending, return_when=asyncio.FIRST_COMPLETED)
            for done_task in done:
                provider = task_to_provider[done_task]
                try:
                    response = done_task.result()
                except Exception as exc:
                    print(_color("system", f"  {provider} error: {exc}"))
                    continue

                if response:
                    if files_included and self._is_stalled(response):
                        print(_color("system", f"─── {provider} ── STALLED (asking for files) — auto-nudging... ───"))
                        stalled_providers.append(provider)
                    elif self._needs_followup(response):
                        # Handle follow-up immediately — brings tab to front
                        new_response = await self._handle_followup(
                            provider, self.sessions[provider], response,
                        )
                        if new_response:
                            header = f"─── {provider} ───"
                            print(_color(provider, header))
                            print(new_response)
                            print()
                            self._save_response(provider, new_response, "initial")
                    else:
                        header = f"─── {provider} ───"
                        print(_color(provider, header))
                        print(response)
                        print()
                        self._save_response(provider, response, "initial")
                else:
                    print(_color(provider, f"─── {provider} ── NO RESPONSE ───"))
                    print()

        # Auto-nudge stalled providers (asking for file uploads)
        if stalled_providers:
            nudge_tasks = {
                provider: asyncio.create_task(
                    self._send_message(self.sessions[provider], _STALL_NUDGE)
                )
                for provider in stalled_providers
                if provider in self.sessions
            }
            for provider, task in nudge_tasks.items():
                try:
                    response = await task
                except Exception:
                    response = None

                if response and not self._is_stalled(response):
                    if self._needs_followup(response):
                        new_response = await self._handle_followup(
                            provider, self.sessions[provider], response,
                        )
                        if new_response:
                            header = f"─── {provider} ───"
                            print(_color(provider, header))
                            print(new_response)
                            print()
                            self._save_response(provider, new_response, "initial")
                    else:
                        header = f"─── {provider} (after nudge) ───"
                        print(_color(provider, header))
                        print(response)
                        print()
                        self._save_response(provider, response, "initial")
                elif response:
                    print(_color("system", f"─── {provider} ── still stuck, showing response ───"))
                    print(response)
                    print()
                else:
                    print(_color(provider, f"─── {provider} ── NO RESPONSE after nudge ───"))
                    print()

    async def _share_responses(self):
        """Cross-pollinate: show each provider what the others said."""
        responses = {
            name: session.last_response
            for name, session in self.sessions.items()
            if session.last_response
        }
        if len(responses) < 2:
            print(_color("system", "Need at least 2 responses to share."))
            return

        print(_color("system", "Sharing responses cross-provider..."))
        print()

        tasks = {}
        for name, session in self.sessions.items():
            others = "\n\n".join(
                f"=== {p} ===\n{r}"
                for p, r in responses.items()
                if p != name
            )
            cross_msg = (
                f"Here's what other AI assistants said about the same question:\n\n"
                f"{others}\n\n"
                "Considering these perspectives, revise your answer. "
                "Where do you agree? Where do you disagree and why? "
                "Give your revised, best answer."
            )
            tasks[name] = asyncio.create_task(self._send_message(session, cross_msg))

        for provider, task in tasks.items():
            response = await task
            if response:
                header = f"─── {provider} (revised) ───"
                print(_color(provider, header))
                print(response)
                print()

    async def _run_consensus(self):
        """Use free-tier API to judge all browser responses."""
        responses = {
            name: session.last_response
            for name, session in self.sessions.items()
            if session.last_response
        }
        if len(responses) < 2:
            print(_color("system", "Need at least 2 responses for consensus."))
            return

        if not self._judge_client:
            self._judge_client, self._judge_model, self._judge_name = _init_judge_client()

        if not self._judge_client:
            print(_color("system", "No free API key found for judging. Set GOOGLE_API_KEY (free) in .env"))
            # Fallback: ask each browser provider to judge
            await self._browser_consensus(responses)
            return

        print(_color("judge", f"Judging via {self._judge_name} API (free tier)..."))
        print()

        responses_text = "\n\n".join(
            f"=== {name} ===\n{resp[:2000]}"
            for name, resp in responses.items()
        )

        judge_prompt = (
            "You are an expert judge comparing AI responses to the same question. "
            "Rate each response for accuracy, completeness, and clarity.\n\n"
            f"Responses:\n{responses_text}\n\n"
            f"Provider names: {list(responses.keys())}\n\n"
            "Respond with EXACTLY this JSON:\n"
            '{"ranking": ["best_provider", "second", ...], '
            '"scores": {"provider": 1-10, ...}, '
            '"synthesis": "Combined best answer from all responses", '
            '"agreements": "What all providers agree on", '
            '"disagreements": "Where they differ and who is more likely correct"}'
        )

        try:
            resp = self._judge_client.chat.completions.create(
                model=self._judge_model,
                messages=[
                    {"role": "system", "content": "You are an expert evaluator comparing AI model responses."},
                    {"role": "user", "content": judge_prompt},
                ],
                max_tokens=2048,
                temperature=0.2,
                timeout=60,
            )
            text = (resp.choices[0].message.content or "").strip()

            # Parse JSON
            if "```" in text:
                for block in text.split("```"):
                    block = block.strip()
                    if block.startswith("json"):
                        block = block[4:].strip()
                    if block.startswith("{"):
                        text = block
                        break

            try:
                result = json.loads(text)
            except json.JSONDecodeError:
                # Show raw judge response
                print(_color("judge", f"─── JUDGE ({self._judge_name}) ───"))
                print(text)
                print()
                return

            # Display results
            print(_color("judge", f"─── SCORES (judged by {self._judge_name}) ───"))
            scores = result.get("scores", {})
            ranking = result.get("ranking", [])
            for i, name in enumerate(ranking):
                score = scores.get(name, "?")
                medal = ["  ", "  ", "  "][i] if i < 3 else "   "
                print(f"  {medal} {_color(name, name)}: {score}/10")
            print()

            if result.get("agreements"):
                print(f"  {_BOLD}Agreements:{_RESET} {result['agreements']}")
            if result.get("disagreements"):
                print(f"  {_BOLD}Disagreements:{_RESET} {result['disagreements']}")
            print()

            if result.get("synthesis"):
                print(_color("consensus", f"─── CONSENSUS (synthesized by {self._judge_name}) ───"))
                print(result["synthesis"])
                print()

            # Record in AttestDB
            if self.db and ranking:
                self._record_consensus_winner(ranking[0], scores.get(ranking[0], 5) / 10)

        except Exception as exc:
            print(_color("system", f"Judging failed: {exc}"))
            print(_color("system", "Falling back to browser-based consensus..."))
            await self._browser_consensus(responses)

    async def _browser_consensus(self, responses: dict[str, str]):
        """Fallback: ask the first browser provider to judge (no API needed)."""
        first_provider = next(iter(self.sessions))
        session = self.sessions[first_provider]

        responses_text = "\n\n".join(
            f"=== {name} ===\n{resp[:2000]}"
            for name, resp in responses.items()
        )

        judge_prompt = (
            "Multiple AI assistants answered the same question. "
            "Compare their responses. Which is best and why? "
            "Synthesize the best answer combining their strengths.\n\n"
            f"Responses:\n{responses_text}"
        )

        print(_color("system", f"Asking {first_provider} to judge (browser fallback)..."))
        response = await self._send_message(session, judge_prompt)
        if response:
            print(_color("consensus", f"─── CONSENSUS (judged by {first_provider}) ───"))
            print(response)
            print()

    async def _run_tournament(self, max_rounds: int = 5):
        """Tournament: providers iteratively improve and score each other until convergence.

        Each round:
          1. Each provider sees the other providers' latest responses
          2. Each provider revises their answer AND scores all responses (including their own)
          3. If all providers agree on the same winner → consensus. Otherwise → next round.
        """
        responses = {
            name: session.last_response
            for name, session in self.sessions.items()
            if session.last_response
        }
        if len(responses) < 2:
            print(_color("system", "Need at least 2 responses to start a tournament. Send a message first."))
            return

        provider_names = list(responses.keys())
        provider_list_str = ", ".join(provider_names)

        for round_num in range(1, max_rounds + 1):
            print(_color("system", f"─── Tournament Round {round_num}/{max_rounds} ───"))
            print()

            # Build per-provider prompts: show others' responses, ask to revise + score
            tasks = {}
            for name, session in self.sessions.items():
                if name not in responses:
                    continue

                others_text = "\n\n".join(
                    f"=== {p} ===\n{r}"
                    for p, r in responses.items()
                    if p != name
                )

                prompt = (
                    f"We put your question into other LLMs and they said the following:\n\n"
                    f"{others_text}\n\n"
                    f"Now do two things:\n"
                    f"1. Give your REVISED best answer, considering what the other models said. "
                    f"Where do you agree? Where do you still disagree and why?\n"
                    f"2. After your answer, on a new line write SCORES: and rate each response "
                    f"(including your own previous one) from 1-10 for accuracy, completeness, "
                    f"and clarity. Format: SCORES: {provider_list_str} = X, Y, Z\n"
                    f"   Example: SCORES: {', '.join(f'{p} = 8' for p in provider_names)}\n\n"
                    f"Start with your revised answer:"
                )
                tasks[name] = asyncio.create_task(self._send_message(session, prompt))

            # Collect revised responses + scores
            scores_by_judge: dict[str, dict[str, int]] = {}

            for provider, task in tasks.items():
                response = await task
                if not response:
                    continue

                # If provider asks follow-up questions during tournament, handle them
                if self._needs_followup(response) and not any(
                    c in response.upper() for c in ["SCORES:", "SCORE:"]
                ):
                    response = await self._handle_followup(
                        provider, self.sessions[provider], response,
                    )
                    if not response:
                        continue

                # Parse score line from response
                revised, provider_scores = self._parse_tournament_response(
                    response, provider_names,
                )
                responses[provider] = revised
                if provider_scores:
                    scores_by_judge[provider] = provider_scores

                # Print and save revised response
                stage = f"round-{round_num}"
                header = f"─── {provider} (round {round_num}) ───"
                print(_color(provider, header))
                print(revised[:3000] if len(revised) > 3000 else revised)
                if provider_scores:
                    scores_str = "  ".join(f"{p}: {s}/10" for p, s in provider_scores.items())
                    print(_color("system", f"  Scores: {scores_str}"))
                print()
                self._save_response(provider, revised, stage)

            # Check convergence: do all judges agree on the same winner?
            if len(scores_by_judge) >= 2:
                winners = {}
                for judge, scores in scores_by_judge.items():
                    if scores:
                        best = max(scores, key=scores.get)
                        winners[judge] = best

                # Count how many judges picked each winner
                winner_counts: dict[str, int] = {}
                for w in winners.values():
                    winner_counts[w] = winner_counts.get(w, 0) + 1

                n_judges = len(winners)
                top_winner = max(winner_counts, key=winner_counts.get) if winner_counts else None
                top_count = winner_counts.get(top_winner, 0) if top_winner else 0

                # Compute average scores
                avg_scores: dict[str, float] = {}
                for p in provider_names:
                    vals = [s.get(p, 0) for s in scores_by_judge.values() if p in s]
                    avg_scores[p] = sum(vals) / len(vals) if vals else 0

                avg_str = "  ".join(f"{p}: {s:.1f}" for p, s in sorted(avg_scores.items(), key=lambda x: -x[1]))
                print(_color("system", f"  Average scores: {avg_str}"))
                self._save_scores(scores_by_judge, f"round-{round_num}")

                # Unanimous or all agree
                if top_count == n_judges:
                    print()
                    print(_color("consensus", f"─── CONSENSUS: {top_winner} wins (unanimous, round {round_num}) ───"))
                    print(responses[top_winner])
                    print()
                    self._record_consensus_winner(top_winner, top_count / n_judges)
                    self._save_final(top_winner, responses[top_winner], {
                        "round": round_num, "type": "unanimous",
                        "scores": avg_scores, "votes": {j: w for j, w in winners.items()},
                    })
                    return

                # Supermajority (all but one agree) with 3+ judges
                if n_judges >= 3 and top_count >= n_judges - 1:
                    dissenters = [j for j, w in winners.items() if w != top_winner]
                    print()
                    print(_color("consensus", f"─── CONSENSUS: {top_winner} wins ({top_count}/{n_judges} agree, round {round_num}) ───"))
                    if dissenters:
                        print(_color("system", f"  Dissent from: {', '.join(dissenters)}"))
                    print(responses[top_winner])
                    print()
                    self._record_consensus_winner(top_winner, top_count / n_judges)
                    self._save_final(top_winner, responses[top_winner], {
                        "round": round_num, "type": "supermajority",
                        "agreement": f"{top_count}/{n_judges}",
                        "dissenters": dissenters,
                        "scores": avg_scores, "votes": {j: w for j, w in winners.items()},
                    })
                    return

                # No consensus yet — show who voted for whom
                vote_str = ", ".join(f"{j} → {w}" for j, w in winners.items())
                print(_color("system", f"  Votes: {vote_str} (no consensus yet)"))
                print()

        # Max rounds reached without full consensus
        print(_color("system", f"─── Tournament ended after {max_rounds} rounds (no full consensus) ───"))
        if scores_by_judge:
            avg_scores = {}
            for p in provider_names:
                vals = [s.get(p, 0) for s in scores_by_judge.values() if p in s]
                avg_scores[p] = sum(vals) / len(vals) if vals else 0
            best_avg = max(avg_scores, key=avg_scores.get)
            print(_color("system", f"  Best by average score: {best_avg} ({avg_scores[best_avg]:.1f}/10)"))
            print()
            print(_color("consensus", f"─── BEST ANSWER: {best_avg} ───"))
            print(responses[best_avg])
            print()
            self._record_consensus_winner(best_avg, avg_scores[best_avg] / 10)
            self._save_scores(scores_by_judge, f"round-{max_rounds}")
            self._save_final(best_avg, responses[best_avg], {
                "round": max_rounds, "type": "best_average",
                "scores": avg_scores,
            })

    def _parse_tournament_response(
        self, response: str, provider_names: list[str],
    ) -> tuple[str, dict[str, int]]:
        """Extract the revised answer and SCORES: line from a tournament response.

        Returns (revised_answer, {provider: score}).
        """
        scores: dict[str, int] = {}
        lines = response.split("\n")

        # Find the SCORES: line (search from the end)
        scores_line_idx = -1
        for i in range(len(lines) - 1, -1, -1):
            if "SCORES:" in lines[i].upper():
                scores_line_idx = i
                break

        if scores_line_idx >= 0:
            scores_text = lines[scores_line_idx]
            # Remove the SCORES: prefix
            scores_part = scores_text.split(":", 1)[-1].strip()

            # Parse "chatgpt = 8, claude = 9, gemini = 7" style
            for provider in provider_names:
                # Look for "provider = N" or "provider: N"
                import re
                pattern = re.compile(
                    rf"{re.escape(provider)}\s*[=:]\s*(\d+)",
                    re.IGNORECASE,
                )
                match = pattern.search(scores_part)
                if match:
                    scores[provider] = min(int(match.group(1)), 10)

            # Revised answer is everything before the SCORES line
            revised = "\n".join(lines[:scores_line_idx]).strip()
        else:
            revised = response.strip()

        return revised, scores

    def _record_consensus_winner(self, winner: str, confidence: float):
        """Record which provider won consensus."""
        if not self.db:
            return
        try:
            from attestdb.core.types import ClaimInput
            self.db.ingest(ClaimInput(
                subject=(winner, "llm_provider"),
                predicate=("won_consensus", "provider_performance"),
                object=("consensus_round", "event"),
                provenance={
                    "source_id": "attest_chat_browser",
                    "source_type": "provider_telemetry",
                },
                confidence=confidence,
                payload={
                    "schema_ref": "consensus_winner/v1",
                    "data": {"confidence": confidence, "mode": "browser"},
                },
            ))
        except Exception:
            pass

    async def _new_chat(self):
        """Navigate all providers to a fresh/new chat page."""
        new_chat_selectors = {
            "chatgpt": [
                'a[data-testid="create-new-chat-button"]',
                'a[href="/"]',
                'nav a[class*="new"]',
            ],
            "claude": [
                'a[href="/new"]',
                'a[data-testid="new-chat"]',
                'button[aria-label="New chat"]',
            ],
            "gemini": [
                'a[href="/app"]',
                'button[aria-label="New chat"]',
                'a[class*="new-chat"]',
            ],
        }

        async def _reset_provider(provider: str, session: BrowserSession):
            page = session.page
            # Try clicking "new chat" button first
            for sel in new_chat_selectors.get(provider, []):
                try:
                    btn = await page.query_selector(sel)
                    if btn:
                        await btn.click()
                        await asyncio.sleep(2)
                        session.last_response = ""
                        session.messages.clear()
                        return True
                except Exception:
                    continue
            # Fallback: navigate directly to the provider's URL
            try:
                url = session.selectors["url"]
                await page.goto(url, wait_until="domcontentloaded", timeout=15000)
                await asyncio.sleep(2)
                session.last_response = ""
                session.messages.clear()
                return True
            except Exception as exc:
                print(_color("system", f"  {provider}: could not start new chat: {exc}"))
                return False

        tasks = [_reset_provider(p, s) for p, s in self.sessions.items()]
        results = await asyncio.gather(*tasks, return_exceptions=True)
        success = sum(1 for r in results if r is True)
        print(_color("system", f"  New chat started on {success}/{len(self.sessions)} providers"))

        # Reset output directory for the new topic
        self._init_output_dir()

    async def _send_to_one(self, provider: str, message: str):
        """Send a message to a single provider and show the response."""
        if provider not in self.sessions:
            print(_color("system", f"  Provider '{provider}' not active. Active: {', '.join(self.sessions.keys())}"))
            return

        session = self.sessions[provider]
        print(_color("system", f"  Sending to {provider} only..."))
        response = await self._send_message(session, message, stream_display=True)

        if response:
            if self._needs_followup(response):
                response = await self._handle_followup(provider, session, response)
            if response:
                # Response was already streamed to terminal — just save it
                # Print a completion line
                print(_color("system", f"  ─── {provider}: {len(response)} chars ───"))
                print()
                self._save_response(provider, response, "targeted")
        else:
            print(_color(provider, f"─── {provider} ── NO RESPONSE ───"))

    async def _screenshot_all(self):
        """Take screenshots of all provider tabs and save to output directory."""
        if not self._output_dir:
            self._init_output_dir()

        ss_dir = self._output_dir / "screenshots"
        ss_dir.mkdir(parents=True, exist_ok=True)

        from datetime import datetime
        timestamp = datetime.now().strftime("%H-%M-%S")

        for provider, session in self.sessions.items():
            try:
                path = ss_dir / f"{provider}_{timestamp}.png"
                await session.page.screenshot(path=str(path), full_page=True)
                print(_color("system", f"  {provider}: saved → {path.name}"))
            except Exception as exc:
                print(_color("system", f"  {provider}: screenshot failed: {exc}"))

        print(_color("system", f"  Screenshots saved to: {ss_dir}/"))

    def _export_session(self):
        """Export the full session as a single markdown document."""
        if not self._output_dir:
            print(_color("system", "No session output to export"))
            return

        from datetime import datetime

        lines = [
            f"# Multi-LLM Chat Session",
            f"**Date:** {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            f"**Providers:** {', '.join(self.sessions.keys())}",
            "",
        ]

        if self._active_models:
            models_str = ", ".join(f"{p}={m}" for p, m in self._active_models.items())
            lines.append(f"**Models:** {models_str}")
            lines.append("")

        # Include conversation history from each provider
        for provider, session in self.sessions.items():
            if not session.messages:
                continue
            lines.append(f"## {provider.title()}")
            lines.append("")
            for msg in session.messages:
                role = msg["role"].title()
                content = msg["content"]
                if role == "User" and len(content) > 500:
                    # Truncate long user messages (file includes)
                    content = content[:500] + "\n\n... (truncated)"
                lines.append(f"### {role}")
                lines.append(content)
                lines.append("")

        # Include final consensus if available
        final_dir = self._output_dir / "final"
        if final_dir.exists():
            best_file = final_dir / "best.md"
            if best_file.exists():
                lines.append("## Final Consensus")
                lines.append(best_file.read_text(encoding="utf-8"))
                lines.append("")

            meta_file = final_dir / "metadata.json"
            if meta_file.exists():
                lines.append("### Tournament Metadata")
                lines.append(f"```json\n{meta_file.read_text(encoding='utf-8')}\n```")
                lines.append("")

        export_path = self._output_dir / "session_report.md"
        export_path.write_text("\n".join(lines), encoding="utf-8")
        print(_color("system", f"  Session exported to: {export_path}"))

        # Open in default editor
        try:
            _open_path(export_path)
        except Exception:
            pass

    def _save_session_state(self, path: Path | None = None):
        """Save session state to JSON for later resumption.

        Saves: messages per provider, active models, auto-followup setting,
        output directory, artifacts.
        """
        if not path:
            path = Path(self.cwd) / ".attest_chat_session.json"

        state = {
            "cwd": self.cwd,
            "providers": list(self.sessions.keys()),
            "active_models": self._active_models,
            "auto_followup": self.auto_followup,
            "output_dir": str(self._output_dir) if self._output_dir else None,
            "sessions": {},
        }
        for name, session in self.sessions.items():
            state["sessions"][name] = {
                "messages": session.messages,
                "last_response": session.last_response[:5000] if session.last_response else "",
                "call_count": session.call_count,
                "artifacts": session.artifacts,
                "rate_limit_hits": session._rate_limit_hits,
            }

        path.write_text(json.dumps(state, indent=2, default=str), encoding="utf-8")
        print(_color("system", f"  Session saved to: {path}"))

    def _load_session_state(self, path: Path | None = None) -> bool:
        """Load session state from JSON.

        Restores messages, models, settings. Does NOT restore browser tabs
        (those need to be re-opened). Returns True if loaded successfully.
        """
        if not path:
            path = Path(self.cwd) / ".attest_chat_session.json"

        if not path.exists():
            print(_color("system", f"  No saved session found at: {path}"))
            return False

        try:
            state = json.loads(path.read_text(encoding="utf-8"))
        except Exception as exc:
            print(_color("system", f"  Failed to load session: {exc}"))
            return False

        # Restore settings
        self._active_models = state.get("active_models", {})
        self.auto_followup = state.get("auto_followup", False)

        saved_output = state.get("output_dir")
        if saved_output and Path(saved_output).exists():
            self._output_dir = Path(saved_output)

        # Restore per-provider state
        for name, sess_state in state.get("sessions", {}).items():
            if name in self.sessions:
                session = self.sessions[name]
                session.messages = sess_state.get("messages", [])
                session.last_response = sess_state.get("last_response", "")
                session.call_count = sess_state.get("call_count", 0)
                session.artifacts = sess_state.get("artifacts", [])
                session._rate_limit_hits = sess_state.get("rate_limit_hits", 0)

        n_restored = sum(
            1 for n in state.get("sessions", {}) if n in self.sessions
        )
        total_msgs = sum(
            len(s.get("messages", [])) for s in state.get("sessions", {}).values()
        )
        print(_color("system", f"  Restored {n_restored} provider(s), {total_msgs} messages"))
        return True

    def _handle_command(self, cmd: str) -> str | None:
        """Parse a / command. Returns command name or None for regular message."""
        parts = cmd.split(maxsplit=1)
        return parts[0].lower() if parts else None

    async def run(self):
        """Main interactive loop."""
        await self._start_browser()

        if not self.sessions:
            print("No providers available. Check your browser logins.")
            return

        # Setup readline tab-completion for @ and /
        self._completer = _setup_readline(self.cwd)

        # Init output directory
        self._init_output_dir()

        providers_str = ", ".join(f"{_BOLD}{name}{_RESET}" for name in self.sessions)
        print(f"\n{_BOLD}attest chat{_RESET} (browser mode, {len(self.sessions)} providers: {providers_str})")
        judge_info = f" | Judge: {self._judge_name} API" if self._judge_name else ""
        if not self._judge_name:
            self._judge_client, self._judge_model, self._judge_name = _init_judge_client()
            judge_info = f" | Judge: {self._judge_name} API (free)" if self._judge_name else " | Judge: browser fallback"
        print(f"Working directory: {self.cwd}{judge_info}")
        print(f"Use @<tab> to pick files/providers | /help for commands")
        models_str = ""
        if self._active_models:
            models_str = " | Models: " + ", ".join(f"{p}={m}" for p, m in self._active_models.items())
        auto_str = " | auto-followup ON" if self.auto_followup else ""
        print(f"Commands: /tournament /new /model /web /export /help{models_str}{auto_str}")
        print()

        # Offer to resume previous session
        session_file = Path(self.cwd) / ".attest_chat_session.json"
        if session_file.exists():
            print(_color("system", "  Previous session found. Type /load to resume, or start fresh."))
            print()

        try:
            while True:
                try:
                    user_input = await asyncio.get_event_loop().run_in_executor(
                        None, lambda: input(f"{_BOLD}>{_RESET} ").strip()
                    )
                except EOFError:
                    break

                if not user_input:
                    continue

                # Multiline input: start with """ or ``` to enter multi-line mode
                if user_input in ('"""', "```"):
                    delimiter = user_input
                    print(_color("system", f"  Multi-line mode (type {delimiter} on a new line to send)"))
                    lines = []
                    while True:
                        try:
                            line = await asyncio.get_event_loop().run_in_executor(
                                None, lambda: input(f"{_color('system', '...')} ")
                            )
                        except EOFError:
                            break
                        if line.strip() == delimiter:
                            break
                        lines.append(line)
                    user_input = "\n".join(lines).strip()
                    if not user_input:
                        continue

                if user_input.startswith("/"):
                    cmd = self._handle_command(user_input)

                    if cmd in ("/quit", "/exit", "/q"):
                        break
                    elif cmd == "/share":
                        await self._share_responses()
                    elif cmd == "/consensus":
                        await self._run_consensus()
                    elif cmd == "/tournament":
                        # Parse optional max_rounds: /tournament 5
                        parts = user_input.split()
                        max_rounds = int(parts[1]) if len(parts) > 1 and parts[1].isdigit() else 5
                        await self._run_tournament(max_rounds=max_rounds)
                    elif cmd == "/providers":
                        for name, session in self.sessions.items():
                            status = "logged in" if session.logged_in else "not logged in"
                            print(f"  {_color(name, name)}: {status}, {session.call_count} messages")
                    elif cmd == "/output":
                        if self._output_dir:
                            print(_color("system", f"Output directory: {self._output_dir}"))
                            _open_path(self._output_dir)
                        else:
                            print(_color("system", "No output directory yet"))
                    elif cmd == "/downloads":
                        if self._output_dir:
                            dl_dir = self._output_dir / "downloads"
                            if dl_dir.exists() and any(dl_dir.iterdir()):
                                print(_color("system", f"Downloads directory: {dl_dir}"))
                                _open_path(dl_dir)
                            else:
                                print(_color("system", "No downloads yet"))
                        else:
                            print(_color("system", "No output directory yet"))
                    elif cmd == "/upload":
                        # Manual file upload: /upload path/to/file.png
                        parts = user_input.split(maxsplit=1)
                        if len(parts) < 2:
                            print(_color("system", "Usage: /upload <file_path> [file_path ...]"))
                        else:
                            paths = []
                            for p in parts[1].split():
                                full = os.path.join(self.cwd, p) if not os.path.isabs(p) else p
                                if os.path.isfile(full):
                                    paths.append(os.path.abspath(full))
                                else:
                                    print(_color("system", f"  File not found: {p}"))
                            if paths:
                                for provider, session in self.sessions.items():
                                    await self._upload_files(session, paths)
                    elif cmd == "/model":
                        # /model chatgpt=o3 claude=opus gemini=pro
                        # /model (no args) — show current models and available options
                        parts = user_input.split()[1:]
                        if not parts:
                            # Show current state and options
                            for p in self.sessions:
                                current = self._active_models.get(p, "default")
                                available = PROVIDER_MODELS.get(p, {})
                                opts = ", ".join(available.keys()) if available else "none"
                                print(f"  {_color(p, p)}: {current} (available: {opts})")
                        else:
                            for spec in parts:
                                if "=" in spec:
                                    prov, model = spec.split("=", 1)
                                    prov = prov.strip().lower()
                                    model = model.strip().lower()
                                    if prov in self.sessions:
                                        await self._select_model(self.sessions[prov], model)
                                    else:
                                        print(_color("system", f"  Provider '{prov}' not active"))
                                else:
                                    print(_color("system", f"  Format: /model provider=model (e.g., /model claude=opus)"))
                    elif cmd == "/new":
                        await self._new_chat()
                    elif cmd == "/screenshot":
                        await self._screenshot_all()
                    elif cmd == "/export":
                        self._export_session()
                    elif cmd == "/web":
                        # Send with web search instruction
                        parts = user_input.split(maxsplit=1)
                        if len(parts) < 2:
                            print(_color("system", "Usage: /web <question>"))
                        else:
                            query = parts[1]
                            web_msg = (
                                "IMPORTANT: Please search the web / use your internet access "
                                "to find the most current information before answering.\n\n"
                                f"{query}"
                            )
                            await self._send_to_all(web_msg)
                    elif cmd == "/auto":
                        # Toggle auto-followup mode
                        self.auto_followup = not self.auto_followup
                        state = "ON" if self.auto_followup else "OFF"
                        print(_color("system", f"  Auto-followup: {state}"))
                        if self.auto_followup:
                            print(_color("system", "  Provider questions will be routed to other providers first"))
                    elif cmd == "/save":
                        parts = user_input.split(maxsplit=1)
                        path = Path(parts[1]) if len(parts) > 1 else None
                        self._save_session_state(path)
                    elif cmd == "/load":
                        parts = user_input.split(maxsplit=1)
                        path = Path(parts[1]) if len(parts) > 1 else None
                        if self._load_session_state(path):
                            print(_color("system", "  Session restored (message history loaded, send a message to continue)"))
                        else:
                            print(_color("system", "  No saved session found"))
                    elif cmd == "/artifacts":
                        total = 0
                        for name, session in self.sessions.items():
                            if session.artifacts:
                                print(_color(name, f"  {name}: {len(session.artifacts)} artifact(s)"))
                                for i, art in enumerate(session.artifacts):
                                    lang = art.get("language", "text")
                                    preview = art.get("content", "")[:80].replace("\n", " ")
                                    print(f"    [{i+1}] {lang}: {preview}...")
                                total += len(session.artifacts)
                            else:
                                print(_color(name, f"  {name}: no artifacts"))
                        if total:
                            print(_color("system", f"  Total: {total} artifact(s). Saved in {self._output_dir}/artifacts/"))
                    elif cmd == "/help":
                        print("Commands:")
                        print("  /share       — One round: share responses cross-provider for revision")
                        print("  /tournament  — Loop: providers improve + score each other until they agree")
                        print("  /tournament N — Same, with max N rounds (default: 5)")
                        print("  /consensus   — Quick judge via free API (Gemini/Groq)")
                        print("  /new         — Start fresh conversations on all providers")
                        print("  /model       — Show/set model tiers (e.g., /model claude=opus chatgpt=o3)")
                        print("  /auto        — Toggle auto-followup (route questions to other providers)")
                        print("  /web <query> — Send with web search instruction")
                        print("  /screenshot  — Screenshot all provider tabs")
                        print("  /export      — Export session as markdown report")
                        print("  /upload FILE — Upload file(s) to all providers")
                        print("  /output      — Open output folder in Finder")
                        print("  /downloads   — Open downloads folder in Finder")
                        print("  /save [path] — Save session state (default: .attest_chat_session.json)")
                        print("  /load [path] — Load a saved session")
                        print("  /artifacts   — List extracted code artifacts")
                        print("  /providers   — List active providers")
                        print("  /quit        — Exit")
                        print()
                        print("Selective sending:")
                        print("  @claude <msg>  — Send to Claude only")
                        print("  @chatgpt <msg> — Send to ChatGPT only")
                        print("  @gemini <msg>  — Send to Gemini only")
                        print()
                        print("Multi-line input:")
                        print('  Type \"\"\" or ``` to start, same delimiter on new line to send')
                    else:
                        print(_color("system", f"Unknown command: {cmd}. Type /help"))
                    continue

                # Check for @provider selective sending (only at start of message)
                selective_match = re.match(
                    r"^@(chatgpt|claude|gemini)\s+(.+)", user_input, re.DOTALL | re.IGNORECASE,
                )
                if selective_match:
                    target_provider = selective_match.group(1).lower()
                    target_msg = selective_match.group(2)
                    resolved = _resolve_file_refs(target_msg, self.cwd)
                    resolved = _auto_include_files(resolved, self.cwd)
                    await self._send_to_one(target_provider, resolved)
                    continue

                # Collect files that need direct upload (images, PDFs, etc.)
                upload_files = _collect_upload_files(user_input, self.cwd)
                if upload_files:
                    names = [os.path.basename(f) for f in upload_files]
                    print(_color("system", f"  (uploading {len(upload_files)} file(s): {', '.join(names)})"))

                # Resolve @file references, then auto-include if user mentions files
                resolved = _resolve_file_refs(user_input, self.cwd)
                resolved = _auto_include_files(resolved, self.cwd)

                # Track if files were included (for stall detection)
                files_included = resolved != user_input
                if files_included:
                    n_files = resolved.count("--- end ")
                    if n_files:
                        print(_color("system", f"  (included {n_files} file(s) from working directory)"))

                # Send to all providers (with file uploads if any)
                await self._send_to_all(resolved, files_included=files_included, upload_files=upload_files or None)

                # Offer tournament if we got 2+ responses
                active_responses = sum(
                    1 for s in self.sessions.values() if s.last_response
                )
                if active_responses >= 2:
                    print(_color("system", "─── All providers responded. ───"))
                    print(_color("system", "  Press Enter to start tournament (providers improve + score each other)"))
                    print(_color("system", "  Or type a message to send to all providers instead"))
                    print()
                    next_input = await asyncio.get_event_loop().run_in_executor(
                        None, lambda: input(f"{_BOLD}tournament or message>{_RESET} ").strip()
                    )
                    if not next_input:
                        # Enter = start tournament
                        await self._run_tournament()
                    elif next_input.startswith("/"):
                        # Handle as command
                        cmd = self._handle_command(next_input)
                        if cmd == "/tournament":
                            parts = next_input.split()
                            max_rounds = int(parts[1]) if len(parts) > 1 and parts[1].isdigit() else 5
                            await self._run_tournament(max_rounds=max_rounds)
                    else:
                        # Send as new message to all
                        resolved = _resolve_file_refs(next_input, self.cwd)
                        resolved = _auto_include_files(resolved, self.cwd)
                        await self._send_to_all(resolved)

        except KeyboardInterrupt:
            print(f"\n{_color('system', 'Shutting down...')}")
        finally:
            # Auto-save session on exit
            has_messages = any(s.messages for s in self.sessions.values())
            if has_messages:
                try:
                    self._save_session_state()
                    print(_color("system", "Session auto-saved. Resume with: attest chat then /load"))
                except Exception as exc:
                    print(_color("system", f"  Warning: failed to auto-save session: {exc}"))
            try:
                if self._output_dir and self._output_dir.exists() and any(self._output_dir.iterdir()):
                    print(_color("system", f"Outputs saved to: {self._output_dir}/"))
                    _open_path(self._output_dir)
            except Exception:
                pass
            await self._stop()

    async def _stop(self):
        """Gracefully close all browser tabs and Playwright."""
        try:
            # Close each page individually first
            for name, session in self.sessions.items():
                try:
                    await session.page.close()
                except Exception:
                    pass

            if self._context:
                try:
                    await self._context.close()
                except Exception:
                    pass

            if self._playwright:
                try:
                    await self._playwright.stop()
                except Exception:
                    pass

            print(_color("system", "Browser closed cleanly."))
        except Exception:
            print(_color("system", "Browser closed."))


def run_browser_chat(
    db=None,
    providers: list[str] | None = None,
    cwd: str | None = None,
    headless: bool = False,
    models: dict[str, str] | None = None,
    auto_followup: bool = False,
):
    """Entry point for browser-based chat. Called from __main__.py."""
    chat = BrowserChat(
        db=db, providers=providers, cwd=cwd, headless=headless,
        models=models, auto_followup=auto_followup,
    )
    try:
        asyncio.run(chat.run())
    except KeyboardInterrupt:
        # asyncio.run can raise KeyboardInterrupt if Ctrl+C during event loop
        print("\nShutdown complete.")
